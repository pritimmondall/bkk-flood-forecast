import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(15, 20, 25)        # #0F1419 Dark Navy
    COLOR_PANEL = RGBColor(26, 32, 39)    # #1A2027 Card BG
    COLOR_BORDER = RGBColor(42, 51, 61)   # #2A333D Border
    COLOR_ACCENT = RGBColor(74, 158, 255)  # #4A9EFF Electric Blue
    COLOR_CORAL = RGBColor(255, 107, 53)  # #FF6B35 Warning Coral
    COLOR_GREEN = RGBColor(58, 183, 149)  # #3AB795 Mint Green
    COLOR_TEXT = RGBColor(230, 237, 243)   # #E6EDF3 Off White
    COLOR_MUTED = RGBColor(139, 152, 165) # #8B98A5 Slate Muted

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text="BANGKOK FLOOD FORECAST V3.0"):
        tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.3))
        tf_cat = tx_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT

        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11), Inches(0.6))
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT

    def add_card(slide, left, top, width, height, bg_color=COLOR_PANEL, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.12), Inches(3.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    tb = slide1.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "BANGKOK FLOOD FORECAST SYSTEM"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT
    p0.space_after = Pt(10)

    p1 = tf.add_paragraph()
    p1.text = "Project Progress & Delivery Roadmap"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT
    p1.space_after = Pt(15)

    p2 = tf.add_paragraph()
    p2.text = "ML-Powered Road Flood Forecasting, Live API Integration & Target Completion Schedule"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_MUTED
    p2.space_after = Pt(30)

    p3 = tf.add_paragraph()
    p3.text = "Target Final Completion: August 24–28, 2026  |  Current Status: Phase 10 (Live Pipeline Completed)"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_GREEN

    # ==========================================
    # SLIDE 2: Executive Summary
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Executive Summary: Project Status & Target Readiness")

    cards_data = [
        ("Core ML & Replay Ready", "53.3% Event POD", "Trained LightGBM onset models achieve 53.3% event POD on 7 years of BMA historical data (2019-2025), accurately predicting road flooding 1-6h ahead.", COLOR_ACCENT),
        ("Live Pipeline Deployed", "Dual-Mode Serving", "Built real-time data ingestion for ThaiWater, Open-Meteo GFS, and Traffy Fondue. Served side-by-side with Replay mode in backend API & dashboard.", COLOR_CORAL),
        ("Final Completion Window", "August 24–28, 2026", "On schedule for final model sign-off, high-resolution GFS downscaling, BMA live integration testing, and system handover.", COLOR_GREEN),
    ]

    for i, (title, highlight, desc, accent) in enumerate(cards_data):
        left = 0.8 + i * 3.9
        add_card(slide2, left, 1.5, 3.7, 5.2)
        
        tb = slide2.shapes.add_textbox(Inches(left + 0.2), Inches(1.7), Inches(3.3), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(15)

        p = tf.add_paragraph()
        p.text = highlight
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = accent
        p.space_after = Pt(15)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 3: Key Accomplishments & Progress
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Major Accomplishments Completed To Date")

    quads = [
        ("1. Data Ingestion & Quality", "Phases 0 - 2", [
            "Ingested 7 years (2019-2025) of BMA raw sensor data.",
            "Covering 131 rain gauges, 300 canal sensors, 107 road sensors.",
            "Engineered quality scorecards & defined 837 true flood events."
        ]),
        ("2. Spatial & External Modeling", "Phases 3 - 4", [
            "Processed 1m DTM elevation, depression depth & TWI.",
            "Integrated Open-Meteo GFS forecast rain (13km grid) & ERA5.",
            "Engineered 50 model-ready features with strict zero-leakage."
        ]),
        ("3. LightGBM ML & Calibration", "Phases 5 & 8", [
            "Trained LightGBM onset specialists across 1, 3, 6h horizons.",
            "Achieved 53.3% Event POD on test year 2025 (Replay mode).",
            "Calibrated probabilities & enforced strict CAP alert rules."
        ]),
        ("4. Live Dual-Mode Architecture", "Phase 10 (Recent)", [
            "Built src/bkkflood/live.py feature assembly pipeline.",
            "Integrated live ThaiWater, GFS forecast & Traffy Fondue feeds.",
            "Deployed live/replay mode toggle in API & React UI."
        ]),
    ]

    for i, (title, tag, points) in enumerate(quads):
        row, col = divmod(i, 2)
        left = 0.8 + col * 5.9
        top = 1.4 + row * 2.7
        add_card(slide3, left, top, 5.7, 2.5)

        tb = slide3.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(5.3), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{title}  |  {tag}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p.space_after = Pt(10)

        for pt in points:
            p = tf.add_paragraph()
            p.text = f"•  {pt}"
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(4)

    # ==========================================
    # SLIDE 4: Live Forecasting vs Replay Analysis
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Live Forecast Capabilities & Input Requirements")

    add_card(slide4, 0.8, 1.4, 6.0, 5.3)
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.6), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Measured Performance Comparison"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(15)

    comp_text = (
        "• Historical Replay Mode: 53.3% Event POD\n"
        "  Uses 7-year archive of BMA 131 rain gauges & road sensors.\n\n"
        "• Current Public Live Mode: 4.9% Event POD\n"
        "  Uses public ThaiWater canal + Open-Meteo GFS 13km rain grid.\n\n"
        "• Key Experimental Insight:\n"
        "  Road sensors are NOT essential (removing them yields 45.1% POD).\n"
        "  The single irreplaceable gap is BMA's 131 Live Rain Gauges."
    )
    p = tf.add_paragraph()
    p.text = comp_text
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED

    add_card(slide4, 7.1, 1.4, 5.4, 5.3)
    tb = slide4.shapes.add_textbox(Inches(7.3), Inches(1.6), Inches(5.0), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Action Plan for Live Precision Boost"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CORAL
    p.space_after = Pt(15)

    action_text = (
        "1. Request Live BMA Rain Gauge Feed\n"
        "   Direct access to BMA's 131 rain gauges (5-min cadence)\n"
        "   will instantly boost live detection from 5% to ~45% POD.\n\n"
        "2. GFS Spatial Downscaling (1-km Grid)\n"
        "   Downscale GFS 13-28km global grid to 1-km resolution to\n"
        "   capture localized 2-5km Bangkok convective storm cells.\n\n"
        "3. Live Sensor Coordinates\n"
        "   Map exact GPS positions for 300 canal sensors to enable\n"
        "   local district-level spatial joins."
    )
    p = tf.add_paragraph()
    p.text = action_text
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT

    # ==========================================
    # SLIDE 5: Project Delivery Timeline (Aug 12 - Aug 28)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Target Completion Timeline: August 24–28, 2026")

    milestones = [
        ("Phase A: Aug 12–15", "GFS Downscaling & API Request", [
            "• Build GFS 1-km spatial downscaling prototype.",
            "• Submit formal request for BMA live rain gauge feed.",
            "• Test high-res rainfall input ingestion."
        ], COLOR_ACCENT),
        ("Phase B: Aug 16–19", "Sequence Model & DB Schema", [
            "• Integrate Phase 6 LSTM sequence model baseline.",
            "• Finalize PostgreSQL / PostGIS database schema.",
            "• Perform full backend load & stress testing."
        ], COLOR_ACCENT),
        ("Phase C: Aug 20–23", "Sandbox & Integration Testing", [
            "• End-to-end sandbox testing with live collectors.",
            "• Verify CAP 1.2 alert generation reliability.",
            "• User acceptance testing & UI dashboard refinement."
        ], COLOR_CORAL),
        ("Phase D: Aug 24–28", "Final Sign-off & Delivery", [
            "• Final model validation & accuracy audit.",
            "• Full documentation & handover package.",
            "• Production deployment & official release."
        ], COLOR_GREEN),
    ]

    for i, (phase, label, tasks, accent) in enumerate(milestones):
        left = 0.8 + i * 2.95
        add_card(slide5, left, 1.5, 2.75, 5.2)

        header_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.1), Inches(1.6), Inches(2.55), Inches(0.7))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = accent
        header_shape.line.fill.background()

        tf = header_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG
        p.alignment = PP_ALIGN.CENTER

        p_sub = tf.add_paragraph()
        p_sub.text = label
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_BG
        p_sub.alignment = PP_ALIGN.CENTER

        tb = slide5.shapes.add_textbox(Inches(left + 0.15), Inches(2.4), Inches(2.45), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        for t_idx, task in enumerate(tasks):
            p = tf.paragraphs[0] if t_idx == 0 else tf.add_paragraph()
            p.text = task
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(8)

    out_path = Path("/Users/pritimmondal/Projects/bkk-flood-forecast/Bangkok_Flood_Forecast_Progress_and_Timeline.pptx")
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == "__main__":
    create_deck()
